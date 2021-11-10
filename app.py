from flask import Flask,jsonify,render_template,request,send_from_directory
import looker_sdk
import os
from pptx import Presentation
from pptx.util import Cm

from docx import Document
from docx.shared import Inches

import PIL.Image as Image
from io import BytesIO


app = Flask(__name__)

sdk=None

@app.route('/download_ppt/<key>',methods = ['GET'])
def download_ppt(key):
    try:
        global sdk
        print(key)
        pptx = Presentation()

        title_only_layout = pptx.slide_layouts[5]
        folder_list=sdk.folder_looks(key)
        if len(folder_list)==0:
            return render_template('no_looks.html')
        for i,look in enumerate(folder_list):
            print(look.id, look.title)
            look_request = {

                    "look_id": look.id, 

                    "result_format": 'png', 

                    "image_width": 960, 

                    "image_height": 540

                }
            
            try:
                image1 = sdk.run_look(**look_request)
                image = Image.open(BytesIO(image1))
                image_file = ''.join([str(look.id), '.png'])
                image.save(image_file) 
            except:
                print(f'Look failed {look.id}: {look.title}')
                image_file = None
            
            pptx.slides.add_slide(title_only_layout)
            pptx.slides[i].shapes.title.text = f''+look.title
            try:
                pptx.slides[i].shapes.add_picture(image_file, Cm(2), Cm(5), width=Cm(20)) # image, left, top, width
                os.remove(image_file)
            except:
                print('Failed to add image to slide')
        pptx.save('data/PPT/demo.pptx')
        return send_from_directory(directory='data/PPT/', path='demo.pptx',as_attachment=True)
    except Exception as e:
        return render_template('error.html',error=str(e))

@app.route('/download_docx/<key>',methods = ['GET'])
def download_docx(key):
    try:
        document = Document()
        document.add_heading("Insights on inventory Items", level=0)
        folder_list=sdk.folder_looks(key)
        if len(folder_list)==0:
            return render_template('no_looks.html')
        for i,look in enumerate(folder_list):
            print(look.id, look.title)
            look_request = {

                "look_id": look.id, 

                "result_format": 'png', 

                "image_width": 960, 

                "image_height": 540

            }
        
            try:
                image1 = sdk.run_look(**look_request)
                image = Image.open(BytesIO(image1))
                image_file = ''.join([str(look.id), '.png'])
                image.save(image_file) 
            except:
                print(f'Look failed {look.id}: {look.title}')
                image_file = None

            document.add_heading(look.title, level=1)
            document.add_picture(image_file, width=Inches(6))
            document.add_page_break()
            try:
                os.remove(image_file)
            except:
                print("error deleting")
        document.save('data/DOCX/'+'demo.docx')
        return send_from_directory(directory='data/DOCX/', path='demo.docx',as_attachment=True)
    except Exception as e:
        return render_template('error.html',error=str(e))



@app.route('/dash',methods = ['POST'])
def folder_page():
    try:
        global sdk
        if request.method == 'POST':
            form_data = request.form.to_dict()
            print(form_data)
            os.environ["LOOKERSDK_BASE_URL"] = str(form_data['base_url'].strip()) #If your looker URL has .cloud in it (hosted on GCP), do not include :19999 (ie: https://your.cloud.looker.com).
            os.environ["LOOKERSDK_API_VERSION"] = "4.0" #3.1 is the default version. You can change this to 4.0 if you want.
            os.environ["LOOKERSDK_VERIFY_SSL"] = "true" #Defaults to true if not set. SSL verification should generally be on unless you have a real good reason not to use it. Valid options: true, y, t, yes, 1.
            os.environ["LOOKERSDK_TIMEOUT"] = "120" #Seconds till request timeout. Standard default is 120.

            #Get the following values from your Users page in the Admin panel of your Looker instance > Users > Your user > Edit API keys. If you know your user id, you can visit https://your.looker.com/admin/users/<your_user_id>/edit.
            os.environ["LOOKERSDK_CLIENT_ID"] =  form_data['client_id'].strip() #No defaults.
            os.environ["LOOKERSDK_CLIENT_SECRET"] = form_data['client_secret'].strip() #No defaults. This should be protected at all costs. Please do not leave it sitting here, even if you don't share this document.
            print("All environment variables set.")
            sdk = looker_sdk.init40()
            print('Looker SDK 4.0 initialized successfully.')
            arr=[]
            for i,items in enumerate(sdk.all_folders()):
                if len(items.name)==0:
                    continue
                temp={"folder_id":items.id,"folder_name":items.name}
                #looks=[ look.title for look in sdk.folder_looks(items.id)]
                #if len(looks)==0:
                #    continue
                #temp["looks_list"]=looks
                print(temp)
                arr.append(temp)

            return render_template('dash.html',data=arr)
    except Exception as e:
        return render_template('error.html',error=str(e))

@app.route('/')
def home_page():
    return render_template('home.html')
	

if __name__ == '__main__':
	app.run()